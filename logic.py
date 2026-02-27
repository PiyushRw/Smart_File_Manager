import os
import pandas as pd
import easyocr
import shutil
import spacy
import torch
from collections import Counter
from sentence_transformers import SentenceTransformer, util
from faster_whisper import WhisperModel
from moviepy.editor import VideoFileClip, AudioFileClip
from markitdown import MarkItDown
from pptx import Presentation
from PIL import Image
import io

# Initialize tools
md = MarkItDown()
ocr = easyocr.Reader(['en'], gpu=False)
nlp = spacy.load("en_core_web_sm", disable=["parser", "ner"])
model = SentenceTransformer('all-MiniLM-L6-v2')
whisper = WhisperModel("base", device="cpu", compute_type="int8")

EXTENSION_MAP = {
    "Images": [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff", ".webp"],
    "Documents": [".pdf", ".docx", ".doc", ".txt", ".md", ".xlsx", ".csv", ".pptx"],
    "Audio": [".mp3", ".wav", ".flac", ".m4a", ".aac"],
    "Video": [".mp4", ".mov", ".mkv", ".avi", ".webm"],
    "Archives": [".zip", ".rar", ".7z", ".tar", ".gz"],
    "Executables": [".exe", ".msi", ".bat", ".sh", ".app"],
    "Code": [".py", ".js", ".html", ".css", ".java", ".cpp"]
}

GENERIC_IGNORE = {"page", "date", "file", "total", "text", "format", "number", "datum", "sheet"}


def _log(message, callback=None):
    """Helper to print or callback"""
    if callback:
        callback(message)
    else:
        print(message)


def extract_images_from_pptx(pptx_path):
    """
    Extract all images from a PowerPoint file and run OCR on them.
    
    Args:
        pptx_path: Path to .pptx file
        
    Returns:
        Combined text from all images in the presentation
    """
    all_image_text = []
    
    try:
        prs = Presentation(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Check if shape contains an image
                if hasattr(shape, "image"):
                    try:
                        # Get image bytes
                        image_bytes = shape.image.blob
                        
                        # Convert to PIL Image
                        image = Image.open(io.BytesIO(image_bytes))
                        
                        # Save temporarily for OCR
                        temp_image_path = f"temp_ppt_image_{slide_num}.png"
                        image.save(temp_image_path)
                        
                        # Run OCR
                        text_list = ocr.readtext(temp_image_path, detail=0)
                        if text_list:
                            all_image_text.extend(text_list)
                        
                        # Clean up temp file
                        if os.path.exists(temp_image_path):
                            os.remove(temp_image_path)
                            
                    except Exception as e:
                        # Skip problematic images
                        continue
                        
    except Exception as e:
        # If can't process PPTX, return empty
        return ""
    
    return " ".join(all_image_text)


def extract_text(file_type, file_path):
    """Extract text from various file types (Documents, Images, Audio, Video).
    
    ENHANCED: Now extracts text from images inside PowerPoint files.
    INCREASED: Word limit raised to 500 words for better semantic matching.
    """
    text = ""
    try:
        # --- A. DOCUMENTS ---
        if file_type == "Documents":
            if file_path.endswith(".txt") or file_path.endswith(".csv"):
                with open(file_path, 'r', encoding='utf-8', errors="ignore") as f:
                    text = f.read(5000)  # Increased from 2000 to 5000 characters
            elif file_path.endswith(".pptx"):
                # Extract text from slides using MarkItDown
                result = md.convert(file_path)
                if result:
                    text = result.text_content
                
                # ENHANCED: Also extract text from images inside the PPTX
                image_text = extract_images_from_pptx(file_path)
                if image_text:
                    text += " " + image_text
            else:
                # Other document types (PDF, DOCX, etc.)
                result = md.convert(file_path)
                if result:
                    text = result.text_content

        # --- B. IMAGES ---
        elif file_type == "Images":
            text_list = ocr.readtext(file_path, detail=0)
            text = " ".join(text_list)

        # --- C. AUDIO & VIDEO ---
        elif file_type == "Audio" or file_type == "Video":
            audio_path = file_path
            temp_audio = "temp_scan.wav"
            CUTOFF_SEC = 150
            created_temp = False

            try:
                clip = None
                if file_type == "Video":
                    clip = VideoFileClip(file_path)
                else:
                    clip = AudioFileClip(file_path)

                if clip.duration:
                    duration_to_read = min(clip.duration, CUTOFF_SEC)
                    sub_clip = clip.subclip(0, duration_to_read)
                    
                    if file_type == "Video":
                        sub_clip.audio.write_audiofile(temp_audio, verbose=False, logger=None)
                    else:
                        sub_clip.write_audiofile(temp_audio, verbose=False, logger=None)
                    
                    sub_clip.close()
                    clip.close()
                    
                    audio_path = temp_audio
                    created_temp = True
                else:
                    clip.close()

            except Exception:
                pass  # Silent fail for trimming
            
            segments, _ = whisper.transcribe(audio_path, beam_size=5)
            text = " ".join([segment.text for segment in segments])

            if created_temp and os.path.exists(temp_audio):
                os.remove(temp_audio)

    except Exception:
        return ""

    # Truncate to 500 words (increased from 200 for better semantic matching)
    if text:
        words = text.split()
        preview = " ".join(words[:500])  # Increased from 200 to 500
        return preview
    
    return ""


def scan_folder(folder_path, progress_callback=None, include_subfolders=True):
    """
    Scan folder and extract metadata + preview text for all files.
    
    Args:
        folder_path: Path to folder to scan
        progress_callback: Optional function(message) for progress updates
        include_subfolders: If True, scan subdirectories; if False, scan only top level
    """
    data = []
    file_count = 0
    
    # Choose scanning method based on include_subfolders
    if include_subfolders:
        # First, count total files (including subfolders)
        total_files = sum([len(files) for _, _, files in os.walk(folder_path)])
        
        # Scan all folders and subfolders
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_count += 1
                file_path = os.path.join(root, file)
                ext = os.path.splitext(file)[1].lower()
                
                # Determine Category
                category = "Others"
                for cat, extensions in EXTENSION_MAP.items():
                    if ext in extensions:
                        category = cat
                        break
                
                # Progress update
                if progress_callback and file_count % 10 == 0:  # Update every 10 files
                    progress_callback(f"⏳ Processing: {file_count}/{total_files} files...")
                
                # Extract Preview Text
                preview_text = ""
                if category in ["Documents", "Images", "Audio", "Video"]:
                    if progress_callback:
                        progress_callback(f"📄 Analyzing: {file}")
                    preview_text = extract_text(category, file_path)
                
                data.append({
                    "Filename": file,
                    "Category": category,
                    "Path": file_path,
                    "Preview": preview_text
                })
    else:
        # Scan only the top-level folder (no subdirectories)
        try:
            files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
            total_files = len(files)
            
            for file in files:
                file_count += 1
                file_path = os.path.join(folder_path, file)
                ext = os.path.splitext(file)[1].lower()
                
                # Determine Category
                category = "Others"
                for cat, extensions in EXTENSION_MAP.items():
                    if ext in extensions:
                        category = cat
                        break
                
                # Progress update
                if progress_callback and file_count % 10 == 0:
                    progress_callback(f"⏳ Processing: {file_count}/{total_files} files...")
                
                # Extract Preview Text
                preview_text = ""
                if category in ["Documents", "Images", "Audio", "Video"]:
                    if progress_callback:
                        progress_callback(f"📄 Analyzing: {file}")
                    preview_text = extract_text(category, file_path)
                
                data.append({
                    "Filename": file,
                    "Category": category,
                    "Path": file_path,
                    "Preview": preview_text
                })
        except Exception as e:
            if progress_callback:
                progress_callback(f"❌ Error scanning folder: {e}")
    
    if progress_callback:
        progress_callback(f"✅ Scanned {file_count} files")
    
    return pd.DataFrame(data)


def extract_keywords_from_preview(df, progress_callback=None):
    """
    Extract top 20 keywords from Preview text for semantic matching.
    Creates a new 'Keywords' column with comma-separated keywords.
    
    INCREASED: From 10 to 20 keywords for better semantic accuracy.
    
    Args:
        df: DataFrame with file data
        progress_callback: Optional function(message) for progress updates
    """
    _log("📝 Extracting keywords from preview text...", progress_callback)
    keywords_list = []
    
    # INCREASED: Top 20 keywords for better matching (was 10)
    TOP_KEYWORDS = 20
    
    for idx, row in df.iterrows():
        if progress_callback and idx % 50 == 0:
            progress_callback(f"🔍 Extracting keywords: {idx}/{len(df)} files...")
            
        preview_text = row['Preview']
        
        if not preview_text or row['Category'] == "Others":
            keywords_list.append("")
            continue
            
        # Process with spaCy
        doc = nlp(preview_text)
        
        # Extract meaningful words
        words = [
            token.lemma_ for token in doc 
            if token.pos_ in ["NOUN", "PROPN", "VERB", "ADJ"]
            and not token.is_stop 
            and not token.is_punct 
            and len(token.text) > 2
            and token.lemma_ not in GENERIC_IGNORE
        ]
        
        # Get top 20 most common keywords (increased from 10)
        if words:
            top_words = [word for word, count in Counter(words).most_common(TOP_KEYWORDS)]
            keywords_list.append(", ".join(top_words))
        else:
            keywords_list.append("")
    
    df['Keywords'] = keywords_list
    _log("✅ Keywords extracted!", progress_callback)
    return df


def refine_categories_with_semantic_search(df, user_query, progress_callback=None):
    """
    Match files to user-specified categories using semantic similarity.
    
    LOGIC:
    1. Extract categories from user query (e.g., "Invoice", "Legal", "Medical")
    2. Extract top 20 keywords from each file's content (increased from 10)
    3. Compare file keywords with query categories using AI embeddings
    4. If similarity > 0.45 → Update category to matched query category
    5. If similarity < 0.45 → Keep original extension-based category
    
    Args:
        df: DataFrame with file data
        user_query: User's categorization query
        progress_callback: Optional function(message) for progress updates
    """
    # Extract target categories from user query
    target_categories = get_categories_from_query(user_query)
    if not target_categories:
        _log("⚠️ No target categories found in query.", progress_callback)
        return df

    _log(f"🎯 Matching files against: {target_categories}", progress_callback)

    # Ensure Keywords column exists
    if 'Keywords' not in df.columns:
        _log("📝 Extracting keywords first...", progress_callback)
        df = extract_keywords_from_preview(df, progress_callback)

    # Encode query categories into AI embeddings
    target_embeddings = model.encode(target_categories, convert_to_tensor=True)
    refined_categories = []

    for idx, row in df.iterrows():
        if progress_callback and idx % 25 == 0:
            progress_callback(f"🔍 Semantic matching: {idx}/{len(df)} files...")
            
        file_keywords = row['Keywords']
        original_category = row['Category']
        
        # Skip if no keywords or already marked as Others
        if not file_keywords or original_category == "Others":
            refined_categories.append(original_category)
            continue

        # Split keywords (top 20 comma-separated words, increased from 10)
        keyword_list = [k.strip() for k in file_keywords.split(",") if k.strip()]
        
        if not keyword_list:
            refined_categories.append(original_category)
            continue
        
        # Encode file keywords into AI embeddings
        keyword_embeddings = model.encode(keyword_list, convert_to_tensor=True)
        
        # Calculate cosine similarity between file keywords and query categories
        cosine_scores = util.cos_sim(keyword_embeddings, target_embeddings)
        max_score = torch.max(cosine_scores).item()

        # DECISION: If strong match (> 0.45) → Use query category
        #           Otherwise → Keep original extension-based category
        if max_score > 0.45:
            best_match_idx = torch.argmax(torch.max(cosine_scores, dim=0).values).item()
            refined_categories.append(target_categories[best_match_idx].capitalize())
        else:
            refined_categories.append(original_category)

    # Update DataFrame with new categories
    df['Category'] = refined_categories
    _log("✅ Semantic refinement complete!", progress_callback)
    return df


def organize_files_into_folders(df, destination_folder, progress_callback=None):
    """
    AUTOMATIC WORKFLOW: Copy files → Verify → Delete originals
    No user choice - this is the only mode of operation.
    
    SAFETY FEATURE: Always copies first, verifies integrity, then deletes originals.
    If any step fails, original files are preserved.
    
    Args:
        df: DataFrame with file data
        destination_folder: Where to organize files
        progress_callback: Optional function(message) for progress updates
    """
    
    if not os.path.exists(destination_folder):
        os.makedirs(destination_folder)
        _log(f"📁 Created main folder: {destination_folder}", progress_callback)

    success_count = 0
    error_count = 0
    total_files = len(df)
    
    # Track which files were successfully copied
    copied_files = []  # List of (source_path, dest_path) tuples
    failed_files = []  # List of failed source paths

    # ===== PHASE 1: COPY ALL FILES =====
    _log("=" * 50, progress_callback)
    _log(f"📋 PHASE 1: Copying {total_files} files to destination...", progress_callback)
    _log("=" * 50, progress_callback)

    for index, row in df.iterrows():
        if progress_callback and (index + 1) % 10 == 0:
            progress_callback(f"📦 Copying: {index + 1}/{total_files} files...")
            
        source_path = row['Path']
        category = row['Category']
        filename = row['Filename']

        if not os.path.exists(source_path):
            _log(f"⚠️ Source file not found: {filename}", progress_callback)
            error_count += 1
            failed_files.append(source_path)
            continue

        category_folder = os.path.join(destination_folder, category)
        
        if not os.path.exists(category_folder):
            os.makedirs(category_folder)
        
        dest_path = os.path.join(category_folder, filename)

        # Handle duplicates
        counter = 1
        name, ext = os.path.splitext(filename)
        while os.path.exists(dest_path):
            dest_path = os.path.join(category_folder, f"{name}_{counter}{ext}")
            counter += 1

        try:
            # Copy file with metadata
            shutil.copy2(source_path, dest_path)
            success_count += 1
            copied_files.append((source_path, dest_path))
            
        except Exception as e:
            _log(f"❌ Error copying {filename}: {e}", progress_callback)
            error_count += 1
            failed_files.append(source_path)

    _log(f"✅ Copy complete: {success_count} files copied, {error_count} errors", progress_callback)

    # ===== PHASE 2: VERIFY COPIED FILES =====
    _log("=" * 50, progress_callback)
    _log(f"🔍 PHASE 2: Verifying {len(copied_files)} copied files...", progress_callback)
    _log("=" * 50, progress_callback)
    
    verification_passed = True
    verified_count = 0
    verification_failed_count = 0
    failed_verifications = []  # Track which files failed verification
    
    for source_path, dest_path in copied_files:
        try:
            # Check if destination file exists
            if not os.path.exists(dest_path):
                _log(f"❌ Verification failed: {os.path.basename(dest_path)} not found", progress_callback)
                verification_failed_count += 1
                verification_passed = False
                failed_verifications.append((source_path, dest_path))
                continue
            
            # Check if file sizes match
            source_size = os.path.getsize(source_path)
            dest_size = os.path.getsize(dest_path)
            
            if source_size != dest_size:
                _log(f"❌ Size mismatch: {os.path.basename(source_path)} (source: {source_size}, dest: {dest_size})", progress_callback)
                verification_failed_count += 1
                verification_passed = False
                failed_verifications.append((source_path, dest_path))
                continue
            
            verified_count += 1
            
            if progress_callback and verified_count % 50 == 0:
                progress_callback(f"✓ Verified: {verified_count}/{len(copied_files)} files...")
                
        except Exception as e:
            _log(f"❌ Verification error for {os.path.basename(source_path)}: {e}", progress_callback)
            verification_failed_count += 1
            verification_passed = False
            failed_verifications.append((source_path, dest_path))
    
    _log(f"✅ Verification complete: {verified_count}/{len(copied_files)} files verified", progress_callback)
    
    if verification_failed_count > 0:
        _log(f"⚠️ {verification_failed_count} files failed verification", progress_callback)

    # ===== PHASE 3: DELETE ORIGINAL FILES (only if verification passed) =====
    if verification_passed and error_count == 0:
        _log("=" * 50, progress_callback)
        _log(f"🗑️  PHASE 3: Deleting {len(copied_files)} original files (all verified)...", progress_callback)
        _log("=" * 50, progress_callback)
        
        deleted_count = 0
        delete_failed_count = 0
        
        for source_path, dest_path in copied_files:
            try:
                if os.path.exists(source_path):
                    os.remove(source_path)
                    deleted_count += 1
                    
                    if progress_callback and deleted_count % 50 == 0:
                        progress_callback(f"🗑️  Deleted: {deleted_count}/{len(copied_files)} original files...")
                        
            except Exception as e:
                _log(f"❌ Failed to delete {os.path.basename(source_path)}: {e}", progress_callback)
                delete_failed_count += 1
        
        _log(f"✅ Deleted {deleted_count} original files", progress_callback)
        
        if delete_failed_count > 0:
            _log(f"⚠️ {delete_failed_count} files could not be deleted (but copies are safe)", progress_callback)
    else:
        # Safety mechanism - don't delete if verification failed
        _log("=" * 50, progress_callback)
        _log("⚠️ SAFETY: Original files NOT deleted", progress_callback)
        _log("=" * 50, progress_callback)
        
        if not verification_passed:
            _log(f"❌ Reason: Verification failed for {verification_failed_count} files", progress_callback)
            _log(f"   Action: Your original files are safe!", progress_callback)
            _log(f"   Copied files are in: {destination_folder}", progress_callback)
            _log(f"   Please manually verify and delete originals if needed", progress_callback)
        
        if error_count > 0:
            _log(f"❌ Reason: {error_count} files had copy errors", progress_callback)
            _log(f"   Action: Your original files are safe!", progress_callback)
            _log(f"   Successfully copied: {success_count} files", progress_callback)
            _log(f"   Failed: {error_count} files", progress_callback)

    # ===== FINAL SUMMARY =====
    _log("=" * 50, progress_callback)
    _log(f"✅ ORGANIZATION COMPLETE!", progress_callback)
    _log(f"   Files processed: {total_files}", progress_callback)
    _log(f"   Successfully copied: {success_count}", progress_callback)
    _log(f"   Copy errors: {error_count}", progress_callback)
    _log(f"   Files verified: {verified_count}", progress_callback)
    
    if verification_passed and error_count == 0:
        _log(f"   Original files deleted: Yes ✓", progress_callback)
    else:
        _log(f"   Original files deleted: No (kept for safety)", progress_callback)
    
    _log(f"   Destination: {destination_folder}", progress_callback)
    _log("=" * 50, progress_callback)


def get_categories_from_query(user_query):
    """Extract target nouns from user query."""
    doc = nlp(user_query)
    targets = [token.text for token in doc if token.pos_ in ["NOUN", "PROPN"] and not token.is_stop]
    return targets


def organize_files_smart(folder_path, destination_folder, user_query=None, include_subfolders=True, progress_callback=None):
    """
    Main orchestration function with progress callbacks.
    
    WORKFLOW:
    1. Scan files → Extension-based categories (Documents, Images, Videos, etc.)
    2. IF user provides query → Match files to query categories using semantic search
    3. ELSE → Keep extension-based categories
    4. Organize files: Copy → Verify → Delete originals
    
    Args:
        folder_path: Source folder to scan
        destination_folder: Where to organize files
        user_query: Optional query like "organize by Invoice and Legal"
        include_subfolders: If True, scan subdirectories; if False, scan only top level
        progress_callback: Optional function(message) for progress updates
    """
    _log("=" * 50, progress_callback)
    _log("🚀 SMART FILE ORGANIZER", progress_callback)
    _log("=" * 50, progress_callback)
    
    # STEP 1: Scan folder
    subfolder_msg = "including subfolders" if include_subfolders else "top-level only"
    _log(f"\n📂 Step 1: Scanning folder ({subfolder_msg})...", progress_callback)
    df = scan_folder(folder_path, progress_callback, include_subfolders)
    _log(f"Found {len(df)} files", progress_callback)
    
    if len(df) == 0:
        _log("⚠️ No files found", progress_callback)
        return df
    
    # STEP 2: Category Refinement (ONLY if query provided)
    if user_query:
        _log("\n🎯 Step 2: Matching files to query categories...", progress_callback)
        _log(f"   Query: '{user_query}'", progress_callback)
        df = refine_categories_with_semantic_search(df, user_query, progress_callback)
    else:
        _log("\n📋 Step 2: Using extension-based categories", progress_callback)
        _log("   (No query provided - files will be organized by type)", progress_callback)
    
    # Show category breakdown
    _log("\n📊 CATEGORY BREAKDOWN:", progress_callback)
    category_counts = df['Category'].value_counts()
    for category, count in category_counts.items():
        _log(f"   • {category}: {count} files", progress_callback)
    
    # STEP 3: Organize files (automatic copy-verify-delete)
    _log(f"\n📦 Step 3: Organizing files (Copy → Verify → Delete)...", progress_callback)
    organize_files_into_folders(df, destination_folder, progress_callback)
    
    _log("\n" + "=" * 50, progress_callback)
    _log("✅ ALL DONE!", progress_callback)
    _log("=" * 50, progress_callback)
    
    return df


# Example usage
if __name__ == "__main__":
    pass